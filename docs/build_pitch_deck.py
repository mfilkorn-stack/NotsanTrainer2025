"""NotSanTrainer 2025 – Pitch-Deck für die ÄLRD-Tagung.
Erzeugt docs/NotSanTrainer2025_Pitch.pptx (max. 10 Folien)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BG       = RGBColor(0x06, 0x08, 0x0F)
SURFACE  = RGBColor(0x0D, 0x12, 0x20)
CARD     = RGBColor(0x11, 0x18, 0x27)
BORDER   = RGBColor(0x1A, 0x22, 0x36)
TEXT     = RGBColor(0xF1, 0xF5, 0xF9)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)
DIM      = RGBColor(0x47, 0x55, 0x69)
PRIMARY  = RGBColor(0xEF, 0x44, 0x44)
ORANGE   = RGBColor(0xF9, 0x73, 0x16)
BLUE     = RGBColor(0x60, 0xA5, 0xFA)
GREEN    = RGBColor(0x10, 0xB9, 0x81)
YELLOW   = RGBColor(0xEA, 0xB3, 0x08)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def bg(slide, rgb=BG):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    fill(r, rgb)
    r.shadow.inherit = False
    return r


def card(slide, x, y, w, h, rgb=CARD):
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    r.adjustments[0] = 0.06
    fill(r, rgb)
    r.line.color.rgb = BORDER
    r.line.width = Pt(0.75)
    return r


def text(slide, x, y, w, h, txt, size=18, bold=False, color=TEXT,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = txt
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def bullets(slide, x, y, w, h, items, size=14, color=MUTED, gap=6):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(gap)
        run = p.add_run()
        run.text = "•  " + item
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def accent_bar(slide, x, y, w):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(3))
    fill(r, PRIMARY)
    return r


def label(slide, x, y, txt, color=PRIMARY):
    return text(slide, x, y, Inches(8), Inches(0.3), txt.upper(),
                size=11, bold=True, color=color)


def footer(slide, n):
    text(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
         "NotSanTrainer 2025  ·  Matthias Filkorn  ·  ÄLRD-Tagung",
         size=9, color=DIM)
    text(slide, Inches(12.0), Inches(7.05), Inches(1), Inches(0.3),
         f"{n} / 10", size=9, color=DIM, align=PP_ALIGN.RIGHT)


# ── Folie 1: Titel ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
glow = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-3), Inches(-2),
                          Inches(8), Inches(8))
fill(glow, RGBColor(0x1A, 0x0A, 0x0A))
logo = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.9), Inches(0.9),
                          Inches(0.8), Inches(0.8))
logo.adjustments[0] = 0.25
fill(logo, PRIMARY)
text(s, Inches(1.85), Inches(1.05), Inches(6), Inches(0.5),
     "NotSanTrainer 2025", size=20, bold=True, color=TEXT)
text(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.4),
     "Die SAA/BPR 2025 als", size=44, bold=True, color=TEXT)
text(s, Inches(0.9), Inches(3.55), Inches(11.5), Inches(1.4),
     "Lern- und Trainingsapp.", size=44, bold=True, color=PRIMARY)
text(s, Inches(0.9), Inches(4.85), Inches(11.5), Inches(0.6),
     "Realistisch trainieren – nicht nur lesen.",
     size=22, color=MUTED)
text(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.4),
     "Pitch · ÄLRD-Tagung 2026", size=14, color=DIM)
text(s, Inches(0.9), Inches(6.0), Inches(11.5), Inches(0.4),
     "Matthias Filkorn  ·  Notfallsanitäter DRK Köln  ·  Gesundheitsökonom",
     size=14, color=MUTED)


# ── Folie 2: Gründerstory ────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
label(s, Inches(0.6), Inches(0.5), "Wie es begonnen hat")
text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
     "Die Idee kam aus dem Dienst.", size=34, bold=True)

card(s, Inches(0.6), Inches(2.0), Inches(7.0), Inches(4.7))
text(s, Inches(0.95), Inches(2.25), Inches(6.4), Inches(0.45),
     "MATTHIAS FILKORN", size=11, bold=True, color=PRIMARY)
text(s, Inches(0.95), Inches(2.6), Inches(6.4), Inches(0.5),
     "Notfallsanitäter beim DRK Köln &", size=18, bold=True)
text(s, Inches(0.95), Inches(2.95), Inches(6.4), Inches(0.5),
     "studierter Gesundheitsökonom.", size=18, bold=True)
bullets(s, Inches(0.95), Inches(3.7), Inches(6.4), Inches(2.8), [
    "Auszubildende fragen nach realistischer Prüfungs-Vorbereitung.",
    "Neue Kollegen wollen die SAA 2025 strukturiert anwenden lernen.",
    "Frische NotSan suchen Sicherheit bei invasiven Maßnahmen.",
    "Erfahrene Kollegen wollen ein Update zur neuen SAA/BPR.",
], size=13, gap=8)

card(s, Inches(8.0), Inches(2.0), Inches(4.7), Inches(4.7))
accent_bar(s, Inches(8.0), Inches(2.0), Inches(4.7))
text(s, Inches(8.3), Inches(2.3), Inches(4.2), Inches(0.4),
     "DAS ERGEBNIS", size=11, bold=True, color=PRIMARY)
text(s, Inches(8.3), Inches(2.7), Inches(4.2), Inches(2.0),
     "In mehreren Nachtschichten gemeinsam mit Claude Code (KI) "
     "ist ein erster funktionsfähiger Prototyp entstanden – "
     "aus der Praxis, für die Praxis.",
     size=15, color=MUTED)
text(s, Inches(8.3), Inches(5.5), Inches(4.2), Inches(0.4),
     "0 → 1 als One-Person-Side-Project.",
     size=14, bold=True, color=ORANGE)
footer(s, 2)


# ── Folie 3: Problem ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
label(s, Inches(0.6), Inches(0.5), "Problem")
text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
     "Die SAA/BPR 2025 ist da – aber wo trainiert man sie?",
     size=30, bold=True)

cols = [
    ("PDF statt Praxis",
     "550+ Seiten SAA/BPR – kein Tool, das daraus echtes Training macht.",
     PRIMARY),
    ("Generische Lern-Apps",
     "Decken alte Algorithmen ab, nicht den Stand der 6-Länder-AG 2025.",
     YELLOW),
    ("Schulen heterogen",
     "Jede Schule baut eigenes Material – kein einheitlicher Standard.",
     BLUE),
    ("Fortbildungspflicht",
     "30 h/Jahr nach §5 NotSanG – nachweisbar trainieren ist aufwendig.",
     GREEN),
]
x0 = Inches(0.6); y0 = Inches(2.3); cw = Inches(2.95); ch = Inches(3.6); gap = Inches(0.15)
for i, (h, t, c) in enumerate(cols):
    cx = x0 + (cw + gap) * i
    card(s, cx, y0, cw, ch)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, y0, cw, Pt(3))
    fill(bar, c)
    text(s, cx + Inches(0.25), y0 + Inches(0.4), cw - Inches(0.5),
         Inches(0.5), h, size=15, bold=True, color=TEXT)
    text(s, cx + Inches(0.25), y0 + Inches(1.05), cw - Inches(0.5),
         Inches(2.4), t, size=12, color=MUTED)

text(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.5),
     "→  Der Bedarf an realitätsnahem, SAA-konformem Training ist im "
     "Rettungsdienst sichtbar – im Dienst und in der Ausbildung.",
     size=14, color=ORANGE, bold=True)
footer(s, 3)


# ── Folie 4: Lösung ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
label(s, Inches(0.6), Inches(0.5), "Lösung")
text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
     "Die SAA/BPR 2025 – als App.", size=34, bold=True)
text(s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.5),
     "1:1 nach dem Stand der 6-Länder-Arbeitsgruppe. "
     "Offline. Ohne Account. Ohne Tracking.",
     size=16, color=MUTED)

pillars = [
    ("Nachschlagen", "Medikamente, invasive Maßnahmen, Krankheitsbilder, "
     "EKG-Befunde — alles im Dienst griffbereit.", PRIMARY),
    ("Trainieren", "1.223 Quizfragen, 306 Fallsimulationen mit Timer, "
     "Algorithmus-Trainer für 34 Behandlungspfade.", BLUE),
    ("Prüfen", "Gesamtprüfung 25 Fragen + Fall, Schwachstellen-Analyse, "
     "Lernfortschritt lokal – ideal für Azubis.", GREEN),
]
x0 = Inches(0.6); y0 = Inches(2.6); cw = Inches(4.05); ch = Inches(3.7); gap = Inches(0.15)
for i, (h, t, c) in enumerate(pillars):
    cx = x0 + (cw + gap) * i
    card(s, cx, y0, cw, ch)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, y0, cw, Pt(3))
    fill(bar, c)
    text(s, cx + Inches(0.3), y0 + Inches(0.4), cw - Inches(0.6),
         Inches(0.6), h, size=20, bold=True, color=c)
    text(s, cx + Inches(0.3), y0 + Inches(1.2), cw - Inches(0.6),
         Inches(2.3), t, size=13, color=MUTED)

text(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
     "Eine Plattform – drei Use Cases: Azubi · aktiver NotSan · Schule.",
     size=13, color=ORANGE, bold=True)
footer(s, 4)


# ── Folie 5: Inhalte / Zahlen ────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, SURFACE)
label(s, Inches(0.6), Inches(0.5), "Was drinsteckt")
text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
     "Vollständig nach SAA/BPR 2025.", size=30, bold=True)

stats = [
    ("1.223", "Quizfragen", "8 Kategorien · mit Erläuterungen", PRIMARY),
    ("306",  "Fallsimulationen", "Training + Prüfung · mit Timer", BLUE),
    ("34",   "Behandlungspfade (BPR)", "Algorithmus-Trainer · klickbar", YELLOW),
    ("29",   "Medikamente", "Dosierungen · KI · Applikation", GREEN),
    ("18",   "Invasive Maßnahmen", "Schritt-für-Schritt · Indikationen", BLUE),
    ("26",   "EKG-Befunde", "Mit echten EKG-Streifen", GREEN),
]
x0 = Inches(0.6); y0 = Inches(2.0); cw = Inches(4.05); ch = Inches(2.2)
gap_x = Inches(0.15); gap_y = Inches(0.2)
for i, (n, h, sub, c) in enumerate(stats):
    col = i % 3
    row = i // 3
    cx = x0 + (cw + gap_x) * col
    cy = y0 + (ch + gap_y) * row
    card(s, cx, cy, cw, ch)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy, cw, Pt(3))
    fill(bar, c)
    text(s, cx + Inches(0.3), cy + Inches(0.25), cw - Inches(0.6),
         Inches(0.9), n, size=42, bold=True, color=c)
    text(s, cx + Inches(0.3), cy + Inches(1.15), cw - Inches(0.6),
         Inches(0.4), h, size=14, bold=True, color=TEXT)
    text(s, cx + Inches(0.3), cy + Inches(1.55), cw - Inches(0.6),
         Inches(0.5), sub, size=11, color=MUTED)

text(s, Inches(0.6), Inches(6.95), Inches(12), Inches(0.3),
     "Datenbasis: SAA/BPR 2025 der ÄLRD aus BW, BB, MV, NRW, SN, ST · "
     "Stand Juli 2025.", size=11, color=DIM)
footer(s, 5)


# ── Folie 6: Funktionen / Differenzierung ───────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
label(s, Inches(0.6), Inches(0.5), "Funktionen & Differenzierung")
text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
     "Was den NotSanTrainer einzigartig macht.", size=28, bold=True)

feats = [
    ("100 % SAA/BPR 2025", "Kein eigenes Curriculum – direkt aus der Quelle "
     "der 6-Länder-Arbeitsgruppe.", PRIMARY),
    ("Offline & PWA-fähig", "Auf Handy, Tablet, Laptop installierbar – "
     "funktioniert ohne Netz.", BLUE),
    ("Kein Account / kein Tracking", "Lernfortschritt bleibt lokal – "
     "datenschutzkonform by design.", GREEN),
    ("Algorithmus-Trainer", "BPRs als klickbare Entscheidungsbäume – "
     "trainierbar, nicht nur lesbar.", YELLOW),
    ("EKG-Modul mit echten Streifen", "26 Befunde mit realen "
     "EKG-Aufzeichnungen.", BLUE),
    ("Fallsimulationen mit Timer", "153 Prüfungsfälle unter Zeitdruck – "
     "wie in der staatl. Prüfung.", PRIMARY),
]
x0 = Inches(0.6); y0 = Inches(2.0); cw = Inches(4.05); ch = Inches(2.3)
gap_x = Inches(0.15); gap_y = Inches(0.2)
for i, (h, t, c) in enumerate(feats):
    col = i % 3
    row = i // 3
    cx = x0 + (cw + gap_x) * col
    cy = y0 + (ch + gap_y) * row
    card(s, cx, cy, cw, ch)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, cx + Inches(0.3),
                             cy + Inches(0.35), Inches(0.25), Inches(0.25))
    fill(dot, c)
    text(s, cx + Inches(0.65), cy + Inches(0.32), cw - Inches(1),
         Inches(0.5), h, size=14, bold=True, color=TEXT)
    text(s, cx + Inches(0.3), cy + Inches(1.0), cw - Inches(0.6),
         Inches(1.2), t, size=12, color=MUTED)
footer(s, 6)


# ── Folie 7: Zielgruppen ────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, SURFACE)
label(s, Inches(0.6), Inches(0.5), "Zielgruppen")
text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
     "Vier Gruppen – ein Standard.", size=30, bold=True)

groups = [
    ("NotSan-Azubis", "Prüfungsvorbereitung mit echten SAA-Fällen, "
     "Quiz nach Lernfeldern, Fallsimulationen mit Timer.", PRIMARY),
    ("Aktive Notfallsanitäter", "Pflicht-Update SAA 2025, schnelles "
     "Nachschlagen im Dienst – offline auf jedem Gerät.", BLUE),
    ("Rettungsdienstschulen", "Einheitliches Lernmaterial, "
     "Lernfortschritts-Analyse, Entlastung der Lehrkräfte.", GREEN),
    ("HiOrgs & RD-Träger", "Verbandsweit einheitlicher Wissensstand, "
     "Unterstützung bei der 30 h-Fortbildungspflicht.", YELLOW),
]
x0 = Inches(0.6); y0 = Inches(2.0); cw = Inches(3.05); ch = Inches(4.3)
gap = Inches(0.15)
for i, (h, t, c) in enumerate(groups):
    cx = x0 + (cw + gap) * i
    card(s, cx, y0, cw, ch)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, y0, cw, Pt(3))
    fill(bar, c)
    text(s, cx + Inches(0.25), y0 + Inches(0.4), cw - Inches(0.5),
         Inches(0.9), h, size=15, bold=True, color=c)
    text(s, cx + Inches(0.25), y0 + Inches(1.4), cw - Inches(0.5),
         Inches(2.7), t, size=12, color=MUTED)

text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
     "Pricing: Free-Tier kostenlos · Plus 19,99 € einmalig · "
     "Schul-/HiOrg-Lizenzen auf Anfrage.",
     size=12, color=ORANGE, bold=True)
footer(s, 7)


# ── Folie 8: Roadmap & Erweiterungen ────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
label(s, Inches(0.6), Inches(0.5), "Erweiterungen & Integrationen")
text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
     "Vom Lerntool zur Trainings-Plattform.", size=28, bold=True)

card(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.7))
accent_bar(s, Inches(0.6), Inches(2.0), Inches(6.0))
text(s, Inches(0.9), Inches(2.25), Inches(5.4), Inches(0.4),
     "AUSBAUSTUFE 1 — CRM-MODUL", size=11, bold=True, color=PRIMARY)
text(s, Inches(0.9), Inches(2.65), Inches(5.4), Inches(0.6),
     "Crew Resource Management", size=22, bold=True, color=TEXT)
text(s, Inches(0.9), Inches(3.3), Inches(5.4), Inches(1.0),
     "Bis zu 80 % schwerer Behandlungsfehler haben "
     "Kommunikations- und Teamversagen als Ursache.",
     size=13, color=MUTED)
bullets(s, Inches(0.9), Inches(4.4), Inches(5.4), Inches(2.2), [
    "Szenariobasiert statt Frontalunterricht",
    "Trainingsbar an denselben Fällen wie Fach-Modul",
    "Nachweis pro Mitarbeiter (§5 NotSanG)",
    "Aggregiertes ÄLRD-Dashboard",
], size=12, gap=6)

card(s, Inches(6.85), Inches(2.0), Inches(6.0), Inches(4.7))
accent_bar(s, Inches(6.85), Inches(2.0), Inches(6.0))
text(s, Inches(7.15), Inches(2.25), Inches(5.4), Inches(0.4),
     "AUSBAUSTUFE 2 — INTEGRATIONEN", size=11, bold=True, color=BLUE)
text(s, Inches(7.15), Inches(2.65), Inches(5.4), Inches(0.6),
     "Anschluss an bestehende Systeme", size=22, bold=True, color=TEXT)
bullets(s, Inches(7.15), Inches(3.4), Inches(5.4), Inches(3.2), [
    "Klassenverwaltung für RD-Schulen",
    "Fortbildungs-Nachweis als PDF / Export ins QM",
    "Single-Sign-On für Träger (DRK, ASB, JUH, MHD, BF)",
    "Anbindung an FMS / CIRS-RD-Auswertungen",
    "Update-Pipeline bei SAA-/BPR-Versionswechseln",
    "Optional: API für E-Learning-Plattformen (Moodle, ILIAS)",
], size=12, gap=5)
footer(s, 8)


# ── Folie 9: Nächste Schritte ───────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, SURFACE)
label(s, Inches(0.6), Inches(0.5), "Nächste Schritte")
text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
     "Vom Prototyp zum Standardwerkzeug.", size=30, bold=True)

steps = [
    ("01", "Pilot in 3 ausgewählten Gebieten",
     "Eine Hand voll Wachen / Schulen je Region — strukturiertes "
     "Feedback aus Azubis, aktiven NotSan und Lehrkräften.", PRIMARY),
    ("02", "Validierung aller Fragen & Fälle",
     "Fachliche Review aller 1.223 Quizfragen und 306 Fälle durch "
     "ÄLRD-Fachgruppe — gegen die SAA/BPR 2025.", BLUE),
    ("03", "Skalierung & CRM-Roll-out",
     "Nach Pilot: Roll-out auf weitere Gebiete, Aktivierung des "
     "CRM-Moduls, Anbindung an RD-Träger-Fortbildungssysteme.", GREEN),
]
x0 = Inches(0.6); y0 = Inches(2.0); cw = Inches(4.05); ch = Inches(4.3)
gap = Inches(0.15)
for i, (num, h, t, c) in enumerate(steps):
    cx = x0 + (cw + gap) * i
    card(s, cx, y0, cw, ch)
    text(s, cx + Inches(0.3), y0 + Inches(0.3), cw - Inches(0.6),
         Inches(1.2), num, size=48, bold=True, color=c)
    text(s, cx + Inches(0.3), y0 + Inches(1.6), cw - Inches(0.6),
         Inches(0.9), h, size=16, bold=True, color=TEXT)
    text(s, cx + Inches(0.3), y0 + Inches(2.5), cw - Inches(0.6),
         Inches(1.7), t, size=12, color=MUTED)

text(s, Inches(0.6), Inches(6.55), Inches(12), Inches(0.4),
     "Was wir von der ÄLRD-Tagung mitnehmen wollen: Pilot-Partner, "
     "fachliche Patenschaft, Empfehlung in Verbänden.",
     size=12, color=ORANGE, bold=True)
footer(s, 9)


# ── Folie 10: CTA / Kontakt ─────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
label(s, Inches(0.6), Inches(0.5), "Lass uns reden")
text(s, Inches(0.6), Inches(1.0), Inches(12), Inches(1.2),
     "Aus dem Dienst heraus gebaut –", size=38, bold=True)
text(s, Inches(0.6), Inches(1.85), Inches(12), Inches(1.2),
     "jetzt für die Praxis öffnen.", size=38, bold=True, color=PRIMARY)

card(s, Inches(0.6), Inches(3.5), Inches(6.0), Inches(3.0))
text(s, Inches(0.9), Inches(3.7), Inches(5.4), Inches(0.4),
     "PILOT-ANFRAGE", size=11, bold=True, color=PRIMARY)
text(s, Inches(0.9), Inches(4.1), Inches(5.4), Inches(0.5),
     "Matthias Filkorn", size=18, bold=True)
text(s, Inches(0.9), Inches(4.55), Inches(5.4), Inches(0.4),
     "Notfallsanitäter · Gesundheitsökonom",
     size=13, color=MUTED)
text(s, Inches(0.9), Inches(5.1), Inches(5.4), Inches(0.5),
     "matthias@madformed.de", size=18, bold=True, color=ORANGE)
text(s, Inches(0.9), Inches(5.65), Inches(5.4), Inches(0.4),
     "MadforMed GmbH", size=12, color=MUTED)

card(s, Inches(6.85), Inches(3.5), Inches(6.0), Inches(3.0))
text(s, Inches(7.15), Inches(3.7), Inches(5.4), Inches(0.4),
     "APP & LANDING", size=11, bold=True, color=BLUE)
text(s, Inches(7.15), Inches(4.1), Inches(5.4), Inches(0.5),
     "mfilkorn-stack.github.io/", size=15, bold=True, color=TEXT)
text(s, Inches(7.15), Inches(4.45), Inches(5.4), Inches(0.5),
     "NotsanTrainer2025/", size=15, bold=True, color=TEXT)
bullets(s, Inches(7.15), Inches(5.1), Inches(5.4), Inches(1.8), [
    "PWA installierbar auf Handy / Tablet",
    "Komplett offline · ohne Account",
    "Free-Tier sofort testbar",
], size=12, gap=4)

text(s, Inches(0.6), Inches(6.85), Inches(12), Inches(0.4),
     "Danke. Fragen? — Lasst uns gemeinsam aus dem Prototyp "
     "einen Standard machen.",
     size=14, color=MUTED, align=PP_ALIGN.CENTER)
footer(s, 10)


out = "docs/NotSanTrainer2025_Pitch.pptx"
prs.save(out)
print(f"OK: {out} ({len(prs.slides)} Folien)")


